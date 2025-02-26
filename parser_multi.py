#!/usr/bin/env python3

import sys
import os
from docx import Document
import chardet
import fitz  # PyMuPDF
import csv
import pandas as pd
from pptx import Presentation
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx.shapes.group import GroupShape
from pptx.enum.shapes import MSO_SHAPE_TYPE

#print("All modules are successfully imported.")

def parse_doc(tmp_file, original_filename):
    #print(f"Parsing file: {tmp_file} with original filename: {original_filename}")  # Debugging output

    # Check if file exists
    if not os.path.exists(tmp_file):
        raise ValueError('File does not exist')

    # Check if file is not empty
    if os.path.getsize(tmp_file) == 0:
        raise ValueError('File is empty')

    #print(f"File extension: {os.path.splitext(original_filename)[1]}")  # Debugging output

    if original_filename.endswith('.txt') or original_filename.endswith('.md') or original_filename.endswith('.json') or original_filename.endswith('.xml'):
        #print("File type: Text-based")  # Debugging output
        return parse_txt(tmp_file)
    elif original_filename.endswith('.docx'):
        #print("File type: DOCX")  # Debugging output
        return parse_docx(tmp_file)
    elif original_filename.endswith('.pptx'):
        #print("File type: PPTX")  # Debugging output
        return parse_pptx(tmp_file)
    elif original_filename.endswith('.pdf'):
        #print("File type: PDF")  # Debugging output
        return parse_pdf(tmp_file)
    elif original_filename.endswith('.csv'):
        #print("File type: CSV")  # Debugging output
        return parse_csv(tmp_file)
    elif original_filename.endswith('.xlsx'):
        #print("File type: XLSX")  # Debugging output
        return parse_xlsx(tmp_file)
    else:
        raise ValueError('File type not supported')

def parse_pdf(file):
    #print("Entering parse_pdf")  # Debugging output

    try:
        doc = fitz.open(file)
        text = ""
        for page_num in range(len(doc)):
            #print(f"Processing page {page_num}")  # Debugging output
            page = doc.load_page(page_num)
            text += page.get_text()
    except Exception as e:
        #print(f"Error extracting text from PDF: {str(e)}")  # Debugging output
        return f"Error extracting text from PDF: {str(e)}"

    if text.strip() == "":
        #print("The file returned no content")  # Debugging output
        return "The file returned no content"
    else:
        #print("PDF parsing completed successfully")  # Debugging output
        return text

def parse_docx(file):
    #print("Entering parse_docx")  # Debugging output

    document = Document(file)
    text = ''
    for item in document.iter_inner_content():
        if isinstance(item, Paragraph):            
            text += item.text + '\n'
        elif isinstance(item, Table):
            text += 'Table'
            for row in item.rows:
                for cell in row.cells:
                    text += cell.text + '\t'
                text += '\n'
    #print("DOCX parsing completed successfully")  # Debugging output
    return text

def check_recursively_for_text(this_set_of_shapes, text_run):
    for shape in this_set_of_shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            check_recursively_for_text(shape.shapes, text_run)
        else:
            if shape.has_text_frame:               
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_run.append(run.text)
            elif shape.has_table:
                for row in shape.table.rows:
                    row_text = ''
                    for cell in row.cells:
                        row_text += cell.text + '\t'
                    text_run.append(row_text)
    return text_run

def parse_pptx(file):
    #print("Entering parse_pptx")  # Debugging output

    presentation = Presentation(file)
    text = []
    for slide in presentation.slides:
        text = check_recursively_for_text(slide.shapes, text)
    #print("PPTX parsing completed successfully")  # Debugging output
    return '\n'.join(text)

def parse_txt(file):
    #print("Entering parse_txt")  # Debugging output

    with open(file, 'rb') as f:
        raw_data = f.read()
        result = chardet.detect(raw_data)
        encoding = result['encoding']

    try:
        with open(file, 'r', encoding=encoding) as f:
            contents = f.read()
    except UnicodeDecodeError:
        with open(file, 'r', encoding='ISO-8859-1') as f:
            contents = f.read()

    #print("Text file parsing completed successfully")  # Debugging output
    return contents

def parse_csv(file):
    #print("Entering parse_csv")  # Debugging output
    delimiter = ','

    contents = ''
    try:
        df = pd.read_csv(file)
        contents = df.to_csv(index=False)
    
    except FileNotFoundError:
        print(f"Error: CSV file '{file}' not found.")

    #print("csv file parsing completed successfully")  # Debugging output
    return contents

def parse_xlsx(file):
    #print("Entering parse_xlsx")  # Debugging output   
    text = ''
    try:
        df = pd.read_excel(file)

        text = df.to_string(index=False)
    except FileNotFoundError:
        print(f"Error: XLSX file '{file}' not found.")

    #print("xlsx parsing completed successfully")  # Debugging output
    return text

if __name__ == '__main__':
    #print("Script started")  # Debugging output
    if len(sys.argv) < 3:
        #print("Usage: python3 parser_multi.py <tmp_file> <original_filename>")  # Debugging output
        sys.exit(1)
    #print(f"File path: {sys.argv[1]}, Original filename: {sys.argv[2]}")  # Debugging output
    text = parse_doc(sys.argv[1], sys.argv[2])
    #print("Parsing completed")  # Debugging output
    print(text)


